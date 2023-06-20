import time
import tqdm
from googletrans import Translator
from pptx import Presentation

def pptx_translate(ppt_file_path: str, source_lang: str, target_lang: str):
    """pptxファイルを翻訳して置き換え保存する

    Args:
        ppt_file_path (str): 翻訳するpptxファイルのパス
        source_lang (str): 翻訳前の言語
        target_lang (str): 翻訳後の言語
    """
    # 計測開始
    start_time = time.time()
    
    # googletransのドライバーを起動
    translator = Translator()

    # pptxファイルを翻訳して置き換え保存
    prs = Presentation(ppt_file_path)
    for slide in tqdm.tqdm(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                try:
                    changed=translator.translate(shape.text,
                                                 dest=target_lang,
                                                 src=source_lang).text
                    shape.text = shape.text.replace(shape.text,changed)
                except:
                    pass
    prs.save(ppt_file_path.replace(".pptx", f"_{source_lang}2{target_lang}.pptx"))
    finish_time = time.time()
    print(f"elapsed time:{finish_time-start_time} sec")

def main():
    pptx_translate("your_pptfile_path", "ja", "en")

if __name__ == "__main__":
    main()